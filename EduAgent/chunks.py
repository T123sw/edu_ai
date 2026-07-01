from __future__ import annotations

import hashlib
import json
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple, Literal

from langchain_core.messages import SystemMessage, HumanMessage
from openai import api_key

from o_agent import get_llm_by_type
from chunk_type import Chunk, ChunkMeta
from define import *


__all__ = [
    'build_chunks',
    'summarize_doc',
    'summarize_ppt',
    'summarize_pdf'
]
# -------------------------
# 你来实现：文件名 -> 绝对路径
# -------------------------
def resolve_abs_path(file_name: str) -> str:
    """
    输入：带后缀的文件名（或相对路径）
    输出：绝对路径（文件必须存在）
    """
    if file_name.endswith(".pdf"):
        return str(PDF_DIR / file_name )
    elif file_name.endswith(".ppt") or file_name.endswith(".pptx"):
        return str(PPT_DIR / file_name )
    else:
        return str(DOC_DIR / file_name )


    raise NotImplementedError


# -------------------------
# utils
# -------------------------
def _sha1(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8", errors="ignore")).hexdigest()


def _norm_text(t: str) -> str:
    if not t:
        return ""
    t = t.replace("\r\n", "\n").replace("\r", "\n")
    lines = [ln.strip() for ln in t.split("\n")]
    lines = [ln for ln in lines if ln]
    return "\n".join(lines).strip()


def _chunk_id(doc_id: str, kind: str, index: int, text: str) -> str:
    return _sha1(f"{doc_id}|{kind}|{index}|{_sha1(text)}")


def _group_lines(lines: List[str], lines_per_chunk: int) -> List[Tuple[int, int, str]]:
    """
    把 lines 按 lines_per_chunk 分组，返回 [(start_line, end_line, text), ...] (line从1开始)
    """
    out: List[Tuple[int, int, str]] = []
    n = len(lines)
    for start in range(0, n, lines_per_chunk):
        part = lines[start:start + lines_per_chunk]
        text = _norm_text("\n".join(part))
        if not text:
            continue
        out.append((start + 1, min(n, start + lines_per_chunk), text))
    return out


# -------------------------
# PDF: 按页
# -------------------------
def pdf_chunks_by_page(pdf_path: str) -> List[Chunk]:
    raise RuntimeError("PDF chunking must use the main MinerU import pipeline; local PDF parsing has been removed.")


# -------------------------
# PPTX: 按页(slide)
# -------------------------
def pptx_chunks_by_slide(pptx_path: str, include_notes: bool = True) -> List[Chunk]:
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("缺少依赖 python-pptx：pip install python-pptx") from e

    p = Path(pptx_path).resolve()
    doc_id = _sha1(str(p))
    prs = Presentation(str(p))
    slides = len(prs.slides)

    chunks: List[Chunk] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        # title
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = _norm_text(slide.shapes.title.text)
        except Exception:
            title = ""

        # all visible texts
        texts: List[str] = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    t = _norm_text(shape.text or "")
                    if t:
                        texts.append(t)
            except Exception:
                pass

        # notes
        notes = ""
        if include_notes:
            try:
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = _norm_text(slide.notes_slide.notes_text_frame.text or "")
            except Exception:
                notes = ""

        parts: List[str] = []
        if title:
            parts.append(f"# {title}")
        if texts:
            parts.append("\n\n".join(texts))
        if include_notes and notes:
            parts.append("[NOTES]\n" + notes)

        slide_text = _norm_text("\n\n".join(parts))
        if not slide_text:
            slide_text = "<EMPTY_SLIDE_TEXT>"

        meta: ChunkMeta = {
            "ext": p.suffix.lower(),
            "slide_no": slide_no,
            "slides": slides,
            "title": title,
            "include_notes": include_notes,
        }

        chunks.append({
            "doc_id": doc_id,
            "chunk_id": _chunk_id(doc_id, "ppt_slide", slide_no, slide_text),
            "source_path": str(p),
            "kind": "ppt_slide",
            "index": slide_no,
            "text": slide_text,
            "meta": meta,
        })

    return chunks

# -------------------------
# DOCX: 先抽段落 -> 拆成“行” -> 按行分组
# -------------------------
def docx_chunks_by_lines(docx_path: str, lines_per_chunk: int = 80) -> List[Chunk]:
    try:
        from docx import Document
    except Exception as e:
        raise RuntimeError("缺少依赖 python-docx：pip install python-docx") from e

    p = Path(docx_path).resolve()
    doc_id = _sha1(str(p))
    doc = Document(str(p))

    # 段落 -> 行（这里把每段当作一行，也可以进一步 split('\n')）
    lines: List[str] = []
    for para in doc.paragraphs:
        t = _norm_text(para.text or "")
        if t:
            # 如果段落里本身带换行，拆开
            for ln in t.split("\n"):
                ln = ln.strip()
                if ln:
                    lines.append(ln)

    if not lines:
        lines = ["<EMPTY_DOCX_TEXT>"]

    grouped = _group_lines(lines, lines_per_chunk)

    chunks: List[Chunk] = []
    for idx, (l0, l1, text) in enumerate(grouped, start=1):
        meta: ChunkMeta = {
            "ext": ".docx",
            "line_range": [l0, l1],
            "source_unit": "docx_paragraphs",
        }
        chunks.append({
            "doc_id": doc_id,
            "chunk_id": _chunk_id(doc_id, "text_lines", idx, text),
            "source_path": str(p),
            "kind": "text_lines",
            "index": idx,
            "text": text,
            "meta": meta,
        })

    return chunks


# -------------------------
# TXT / 其他：按行分组
# -------------------------
def text_file_chunks_by_lines(text_path: str, lines_per_chunk: int = 50) -> List[Chunk]:
    p = Path(text_path).resolve()
    doc_id = _sha1(str(p))

    try:
        raw = p.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        raw = p.read_text(errors="ignore")

    lines = raw.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    if not any(ln.strip() for ln in lines):
        lines = ["<EMPTY_TEXT_FILE>"]

    grouped = _group_lines(lines, lines_per_chunk)

    chunks: List[Chunk] = []
    for idx, (l0, l1, text) in enumerate(grouped, start=1):
        meta: ChunkMeta = {
            "ext": p.suffix.lower(),
            "line_range": [l0, l1],
            "source_unit": "file",
        }
        chunks.append({
            "doc_id": doc_id,
            "chunk_id": _chunk_id(doc_id, "text_lines", idx, text),
            "source_path": str(p),
            "kind": "text_lines",
            "index": idx,
            "text": text,
            "meta": meta,
        })

    return chunks


# -------------------------
# 统一入口：按后缀分派
# -------------------------
def build_chunks(
    file_name: str,
    *,
    resolve_path_fn: Callable[[str], str] = resolve_abs_path,
    lines_per_chunk: int = 80,
    include_notes: bool = True,
) -> List[Chunk]:

    abs_path = Path(resolve_path_fn(file_name)).resolve()
    if not abs_path.exists():
        raise FileNotFoundError(str(abs_path))

    ext = abs_path.suffix.lower()

    if ext == ".pdf":
        return pdf_chunks_by_page(str(abs_path))

    if ext in (".pptx", ".ppt"):
        # 简化：.ppt 先不支持直接解析（要转pptx）。你如果需要我再加LibreOffice转。
        if ext == ".ppt":
            raise RuntimeError("暂不直接支持 .ppt，请先转为 .pptx（LibreOffice可批量转）")
        return pptx_chunks_by_slide(str(abs_path), include_notes=include_notes)

    if ext == ".docx":
        return docx_chunks_by_lines(str(abs_path), lines_per_chunk=lines_per_chunk)

    # 其他全部当“文本文件”：按行
    return text_file_chunks_by_lines(str(abs_path), lines_per_chunk=lines_per_chunk)


# -------------------------
# 你写：chunks -> 全文摘要
# -------------------------
def summarize_all_chunks(chunks: List[Chunk]) -> str:
    """
    你来实现：接收 chunks，返回完整摘要（str）
    """
    pass



    raise NotImplementedError


import asyncio, json
from typing import List, Dict, Any
from langchain_core.messages import SystemMessage, HumanMessage


async def _ainvoke(llm, messages):
    """优先用 ainvoke，没有就线程异步化。"""
    if hasattr(llm, "ainvoke"):
        r = await llm.ainvoke(messages)
    else:
        r = await asyncio.to_thread(llm.invoke, messages)
    return getattr(r, "content", r)


async def summarize_ppt(chunks: List[Chunk], file_name: str) -> str:
    from o_agent import get_llm_from_config
    llm = get_llm_from_config(temperature=0.3)

    sem = asyncio.Semaphore(8)
    async def one(chunk: Chunk):
        async with sem:
            msg = [
                SystemMessage(content="你负责帮用户总结ppt页面的信息。"),
                HumanMessage(content=f"当前ppt的页号是{chunk['index']}\n文本内容：\n{chunk['text']}"),
            ]
            content = await _ainvoke(llm, msg)
            return {"index": chunk["index"], "summary": content}

    # 并发汇总
    results = await asyncio.gather(*[one(c) for c in chunks])
    results.sort(key=lambda x: x["index"])
    js = json.dumps(results, ensure_ascii=False)

    final_msg = [
        SystemMessage(content="你负责帮用户总结整个ppt内容。用户提供一个json对象列表，每个对象包含index和summary。请生成总体总结，如：这个[文件名]ppt讲了..."),
        HumanMessage(content=f"文件名：{file_name}\njson对象：{js}"),
    ]
    return await _ainvoke(llm, final_msg)


async def summarize_pdf(chunks: List[Chunk], file_name: str) -> str:
    from o_agent import get_llm_from_config
    llm = get_llm_from_config(temperature=0.3)

    sem = asyncio.Semaphore(8)
    async def one(chunk: Chunk):
        async with sem:
            msg = [
                SystemMessage(content="你负责帮用户总结pdf页面的信息。"),
                HumanMessage(content=f"当前pdf的页号是{chunk['index']}\n文本内容：\n{chunk['text']}"),
            ]
            content = await _ainvoke(llm, msg)
            return {"index": chunk["index"], "summary": content}

    results = await asyncio.gather(*[one(c) for c in chunks])
    results.sort(key=lambda x: x["index"])
    js = json.dumps(results, ensure_ascii=False)

    final_msg = [
        SystemMessage(content="你负责帮用户总结整个pdf内容。用户提供一个json对象列表，每个对象包含index和summary。请生成总体总结，如：这个[文件名]讲了..."),
        HumanMessage(content=f"文件名：{file_name}\njson对象：{js}"),
    ]
    return await _ainvoke(llm, final_msg)


async def summarize_doc(chunks: List[Chunk], file_name: str) -> str:
    from o_agent import get_llm_from_config
    llm = get_llm_from_config(temperature=0.3)

    sem = asyncio.Semaphore(8)
    async def one(chunk: Chunk):
        async with sem:
            lr = chunk.get("meta", {}).get("line_range")
            msg = [
                SystemMessage(content="你负责帮用户总结文本（doc/txt）的内容。"),
                HumanMessage(content=f"当前文本的行数范围：{lr}\n文本内容：\n{chunk['text']}"),
            ]
            content = await _ainvoke(llm, msg)
            return {"index": chunk["index"], "summary": content}

    results = await asyncio.gather(*[one(c) for c in chunks])
    results.sort(key=lambda x: x["index"])
    js = json.dumps(results, ensure_ascii=False)

    final_msg = [
        SystemMessage(content="你负责帮用户总结整个文本文档内容。用户提供一个json对象列表，每个对象包含index和summary。请生成总体总结，如：这个[文件名]讲了..."),
        HumanMessage(content=f"文件名：{file_name}\njson对象：{js}"),
    ]
    return await _ainvoke(llm, final_msg)



